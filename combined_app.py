#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
combined_app.py

目的:
  app.pyとmain.pyを統合したセキュリティ規則チェックRAGチャットボット
  CUIモードとStreamlit WebUIの両方をサポートします。

統合機能:
  - コマンドライン引数でCUIモード（--cui）またはWebUIモード（デフォルト）を選択
  - セキュリティ規則一覧ルールをチェックする AgentRAG
  - ドキュメント要約・確認エージェント（Agent A/B）
  - LCEL（Runnable 等）を利用したチェーン
  - OpenAI の gpt-4o を ChatOpenAI で呼び出し（OPENAI_API_KEY 必須）

使用方法:
  python combined_app.py           # Streamlit WebUI モード
  python combined_app.py --cui     # CUI モード
"""

import os
import sys
import json
import glob
import re
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import textwrap
from datetime import datetime
import argparse

# --- 依存パッケージのインポート（利用環境でインストールされている前提） ---
try:
    # LangChain とモデルラッパー
    from langchain.chat_models import ChatOpenAI
    from langchain.schema import HumanMessage, SystemMessage
    from langchain.embeddings import HuggingFaceEmbeddings
    from langchain.vectorstores import Chroma
    from langchain.text_splitter import CharacterTextSplitter
    from langchain.docstore.document import Document
except Exception as e:
    print("必要なライブラリが見つかりません: langchain 等。\n`pip install -r requirements.txt` を実行してください。\nエラー: ", e)
    sys.exit(1)

# LCEL 系のインポート（利用可能なら利用する）
USE_LCEL = True
import importlib

Runnable = None
RunnablePassthrough = None
try:
    mod = importlib.import_module("langchain_experimental")
    Runnable = getattr(mod, "Runnable", None)
    RunnablePassthrough = getattr(mod, "RunnablePassthrough", None)
except Exception:
    try:
        mod2 = importlib.import_module("langchain.experimental.runnable")
        Runnable = getattr(mod2, "Runnable", None)
        RunnablePassthrough = getattr(mod2, "RunnablePassthrough", None)
    except Exception:
        Runnable = None
        RunnablePassthrough = None

if Runnable is None or RunnablePassthrough is None:
    USE_LCEL = False

import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 環境変数から API キーを取得
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    print("環境変数 OPENAI_API_KEY が設定されていません。設定してから再実行してください。")
    sys.exit(1)

# --- 設定値 ---
BASE_DIR = Path(__file__).parent
RULE_DIR = BASE_DIR / "rule"
SPEC_DIR = BASE_DIR / "specification"
CHROMA_DIR = BASE_DIR / "chroma_db"
CHROMA_COLLECTION = "specs"

# Embedding モデル名（ローカルで実行できる軽量モデルを使用）
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"

# RAG 検索時に取り出すドキュメント数
TOP_K = 3


def load_rules_from_dir(rule_dir: Path) -> List[Dict[str, Any]]:
    """rule ディレクトリ配下の全 JSON を再帰的に読み込み、ルールをフラットなリストで返す。

    ルール JSON は配列 または {"rules": [...] } の形式に対応。
    再帰的な構造（子ルールを 'children' などで持つ）もフラット化して返す。
    
    新機能: touitsukijun_r7.json の階層構造 (sections -> subsections -> items) に対応
    """
    rules: List[Dict[str, Any]] = []
    
    def create_rule_entry(item: Dict[str, Any], parent_path: str = "", file_source: str = "") -> Dict[str, Any]:
        """個別ルールエントリを作成"""
        rid = item.get("id") or item.get("rule_id") or item.get("name") or None
        title = item.get("title") or item.get("name") or rid or "unnamed"
        path_label = f"{parent_path}/{title}" if parent_path else title
        
        # content 構築: description + information (あれば)
        content_parts = []
        if item.get("description"):
            content_parts.append(item["description"])
        if item.get("information"):
            content_parts.append(f"\n[詳細情報]\n{item['information']}")
        if item.get("content"):
            content_parts.append(item["content"])
            
        content = "\n".join(content_parts) if content_parts else json.dumps(item, ensure_ascii=False)
        
        return {
            "id": rid,
            "title": title,
            "path": path_label,
            "content": content,
            "type": item.get("type", "未分類"),
            "source_file": file_source,
            "raw": item,
        }
    
    for path in rule_dir.rglob("*.json"):
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception as e:
            logger.warning(f"ルールファイルを読み込めませんでした: {path} - {e}")
            continue

        file_source = path.name
        
        # 入れ子対応: data が配列か辞書か
        candidates = []
        if isinstance(data, list):
            candidates = data
        elif isinstance(data, dict):
            if "rules" in data and isinstance(data["rules"], list):
                candidates = data["rules"]
            else:
                candidates = [data]

        # 階層構造の再帰的処理
        def walk(item: Dict[str, Any], parent_path: str = ""):
            # 現在の項目がルールとして追加すべきものかチェック
            has_description = bool(item.get("description") or item.get("content"))
            
            if has_description:
                # ルールエントリとして追加
                entry = create_rule_entry(item, parent_path, file_source)
                rules.append(entry)
            
            # 階層を下って子要素を処理
            current_path = f"{parent_path}/{item.get('title', item.get('id', ''))}" if parent_path else (item.get('title') or item.get('id') or "")
            
            # 新しい階層キー: sections, subsections, items に対応
            for child_key in ("children", "rules", "subrules", "items", "sections", "subsections"):
                if child_key in item and isinstance(item[child_key], list):
                    for child in item[child_key]:
                        if isinstance(child, dict):
                            walk(child, current_path)

        for it in candidates:
            if isinstance(it, dict):
                walk(it)

    logger.info(f"読み込んだルール数: {len(rules)}")
    return rules


def text_from_pdf(path: Path) -> str:
    """シンプルな PDF テキスト抽出。pypdf を利用。ページごとに連結する。"""
    try:
        import pypdf
    except Exception:
        raise RuntimeError("pypdf が必要です。pip install pypdf を実行してください。")
    text_parts = []
    try:
        reader = pypdf.PdfReader(str(path))
        for p in reader.pages:
            txt = p.extract_text() or ""
            text_parts.append(txt)
    except Exception as e:
        logger.warning(f"PDF 読み込み失敗 {path}: {e}")
    return "\n".join(text_parts)


def text_from_docx(path: Path) -> str:
    try:
        import docx
    except Exception:
        raise RuntimeError("python-docx が必要です。pip install python-docx を実行してください。")
    try:
        doc = docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        logger.warning(f"DOCX 読み込み失敗 {path}: {e}")
        return ""


def text_from_xlsx(path: Path) -> str:
    """Excel ファイル(.xlsx)からテキストを抽出"""
    try:
        import openpyxl
    except Exception:
        raise RuntimeError("openpyxl が必要です。pip install openpyxl を実行してください。")
    
    text_parts = []
    try:
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"\n--- シート: {sheet_name} ---\n")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if row_text:
                    text_parts.append("\t".join(row_text))
        workbook.close()
    except Exception as e:
        logger.warning(f"XLSX 読み込み失敗 {path}: {e}")
        return ""
    
    return "\n".join(text_parts)


def text_from_pptx(path: Path) -> str:
    """PowerPoint ファイル(.pptx)からテキストを抽出"""
    try:
        from pptx import Presentation
    except Exception:
        raise RuntimeError("python-pptx が必要です。pip install python-pptx を実行してください。")
    
    text_parts = []
    try:
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- スライド {i} ---\n")
            
            # スライド内の全テキストを抽出
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
                
                # 表がある場合のテキスト抽出
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text_parts.append("\t".join(row_text))
    except Exception as e:
        logger.warning(f"PPTX 読み込み失敗 {path}: {e}")
        return ""
    
    return "\n".join(text_parts)


def load_spec_documents(spec_dir: Path) -> List[Document]:
    """`specification/` 配下のドキュメントを読み込み、langchain Document のリストを返す。

    対応: pdf, docx, xlsx, pptx, md, txt
    メモリ節約: ファイル毎にチャンク分割を行い、最低限のメタデータを付与
    """
    docs: List[Document] = []
    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    for path in spec_dir.rglob("*"):
        if path.is_dir():
            continue
        lower = path.suffix.lower()
        try:
            if lower == ".pdf":
                text = text_from_pdf(path)
            elif lower == ".docx":
                text = text_from_docx(path)
            elif lower == ".xlsx":
                text = text_from_xlsx(path)
            elif lower == ".pptx":
                text = text_from_pptx(path)
            elif lower in (".md", ".txt"):
                text = path.read_text(encoding="utf-8", errors="ignore")
            else:
                logger.debug(f"未対応ファイル形式をスキップ: {path}")
                continue
        except Exception as e:
            logger.warning(f"ファイル読み込み失敗 {path}: {e}")
            continue

        if not text.strip():
            logger.debug(f"空のファイルをスキップ: {path}")
            continue

        chunks = text_splitter.split_text(text)
        for i, c in enumerate(chunks):
            meta = {"source": str(path), "chunk": i, "file_type": lower}
            docs.append(Document(page_content=c, metadata=meta))

    logger.info(f"読み込んだドキュメントチャンク数: {len(docs)}")
    return docs


def init_chroma(docs: List[Document]) -> Chroma:
    """ChromaDB を初期化して、ドキュメントを格納する。既存コレクションがあれば再利用。

    埋め込みはローカル HuggingFace モデルを使う（軽量モデル推奨）
    """
    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
    vectordb = Chroma(persist_directory=str(CHROMA_DIR), collection_name=CHROMA_COLLECTION, embedding_function=embeddings)

    try:
        existing = vectordb._collection.count() if hasattr(vectordb, "_collection") else None
    except Exception:
        existing = None

    if existing in (None, 0):
        if docs:
            logger.info("Chroma にドキュメントを追加します...")
            vectordb.add_documents(docs)
            vectordb.persist()
    else:
        logger.info("既存の Chroma コレクションを利用します。新規追加は行いません。")

    return vectordb


def retrieve_related_docs(vectordb: Chroma, query: str, k: int = TOP_K, include_rules: bool = False) -> List[Document]:
    """関連ドキュメントを検索する。"""
    docs = []
    
    try:
        main_docs = vectordb.similarity_search(query, k=k)
        docs.extend(main_docs)
    except Exception as e:
        logger.warning(f"メインドキュメント検索エラー: {e}")
    
    return docs[:k]


def make_chat_model() -> ChatOpenAI:
    """ChatOpenAI を作成。モデル名は gpt-4o を指定する。"""
    llm = ChatOpenAI(model_name="gpt-4o", temperature=0.0)
    return llm


def agent_a_summarize(llm: ChatOpenAI, rule_text: str, docs: List[Document]) -> str:
    """Agent A: ドキュメント要約エージェント"""
    context = "\n\n---関連ドキュメント---\n"
    for d in docs[:TOP_K]:
        src = d.metadata.get("source") if d.metadata else "<unknown>"
        context += f"[source: {src}]\n{d.page_content}\n\n"

    system_prompt = (
        "あなたは優秀なドキュメント要約者です。以下のセキュリティルールを読み、"
        "関連ドキュメントの内容を完全性を保って要約してください。重要な条件、要件、及び検証ポイントを箇条書きで示してください。"
        "\n\n重要: 回答は必ず日本語で行ってください。出力の本文は日本語で記載し、読みやすい箇条書きを心がけてください。"
    )
    messages = [SystemMessage(content=system_prompt), HumanMessage(content=f"ルール:\n{rule_text}\n\nドキュメントコンテキスト:\n{context}")]
    resp = llm(messages)
    return resp.content


def agent_b_check(llm: ChatOpenAI, rule_summary: str, rule_raw: Dict[str, Any], docs: List[Document]) -> Dict[str, Any]:
    """Agent B: ドキュメント確認エージェント"""
    context = "\n\n".join([f"[src:{d.metadata.get('source')}]\n{d.page_content}" for d in docs[:TOP_K]])

    strict_prompt = """
あなたは技術的な評価者です。以下のルール要約と元ルールを読み、与えられたシステム情報がそのルールに従っているかを評価してください。

出力は厳密な JSON のみを返してください。余計な説明や追加テキストは一切書かず、必ず純粋な JSON テキストだけを返してください（コードフェンスや説明を含めないでください）。

JSON スキーマ例:
{
    "result": "〇|△|×",
    "evidence": [
        {"source": "ファイル名や識別子", "excerpt": "抜粋テキスト..."}
    ],
    "details": "追加の説明(任意)"
}

重要: JSON のキー名は英語のままにし、値や説明文は日本語で記載してください。
"""

    init_human = HumanMessage(content=f"ルール要約:\n{rule_summary}\n\n元ルール(raw):\n{json.dumps(rule_raw, ensure_ascii=False)}\n\nドキュメントコンテキスト:\n{context}")

    messages = [SystemMessage(content=strict_prompt), init_human]
    resp = llm(messages)
    text = resp.content

    logger.debug("Agent B raw output (head 1000 chars): %s", text[:1000].replace('\n', '\\n'))

    def _save_model_output(rule_id: str, content: str):
        try:
            logs_dir = BASE_DIR / "logs"
            logs_dir.mkdir(exist_ok=True)
            ts = datetime.utcnow().strftime("%Y%m%dT%H%M%SZ")
            fname = logs_dir / f"agent_b_output_{ts}_{str(rule_id)[:60].replace(' ', '_')}.log"
            with open(fname, "w", encoding="utf-8") as lf:
                lf.write("--- RAW MODEL OUTPUT ---\n")
                lf.write(content)
            logger.info("モデル出力をログに保存しました: %s", fname)
        except Exception as e:
            logger.debug("モデル出力ログ保存に失敗しました: %s", e)

    if os.environ.get("SAVE_MODEL_OUTPUT", "1") != "0":
        try:
            _save_model_output(rule_raw.get("id") or rule_raw.get("title") or "unknown", text)
        except Exception:
            pass

    parsed = None
    try:
        parsed = json.loads(text)
    except Exception:
        parsed = None

    retries = 0
    while parsed is None and retries < 2:
        retries += 1
        logger.info("JSON パース失敗: モデルへ再試行を行います (試行 %d)。", retries)
        followup = (
            "前の回答は有用でしたが、要求された通り厳密な JSON のみで出力されていませんでした。"
            "以下の JSON スキーマに厳密に合わせ、純粋な JSON テキストのみを出力してください。"
            "\n\nスキーマ: {\"result\":\"〇|△|×\", \"evidence\": [ {\"source\":..., \"excerpt\":...} ], \"details\": \"任意の文字列\" }"
            "\n\n元の出力を参照して、上記スキーマにマッピングして JSON のみを返してください。"
        )
        follow_messages = [SystemMessage(content=strict_prompt), init_human, HumanMessage(content=followup + "\n\n前の出力:\n" + text)]
        try:
            resp2 = llm(follow_messages)
            text2 = resp2.content
            logger.debug("Agent B retry raw output (head 1000): %s", text2[:1000].replace('\n', '\\n'))
            try:
                parsed = json.loads(text2)
                text = text2
                break
            except Exception:
                m2 = re.search(r"(\{[\s\S]*\})", text2)
                if m2:
                    try:
                        parsed = json.loads(m2.group(1))
                        text = m2.group(1)
                        break
                    except Exception:
                        parsed = None
                t2 = text2.replace("'", '"')
                t2 = re.sub(r",\s*([}\]])", r"\1", t2)
                t2 = re.sub(r'([\{,\s])(\w+)\s*:', r'\1"\2":', t2)
                try:
                    parsed = json.loads(t2)
                    text = t2
                    break
                except Exception:
                    parsed = None
        except Exception as e:
            logger.debug("モデル再試行中に例外: %s", e)
            parsed = None

    if parsed is None:
        m = re.search(r"(\{[\s\S]*\})", text)
        if m:
            candidate = m.group(1)
            try:
                parsed = json.loads(candidate)
            except Exception:
                parsed = None

        if parsed is None:
            t2 = text.replace("'", '"')
            t2 = re.sub(r",\s*([}\]])", r"\1", t2)
            t2 = re.sub(r'([\{,\s])(\w+)\s*:', r'\1"\2":', t2)
            try:
                parsed = json.loads(t2)
            except Exception:
                parsed = None

    if parsed is None:
        logger.warning("モデルの出力を JSON としてパースできませんでした。ヒューリスティック抽出を試みます。")

        def _heuristic_parse(text: str) -> Dict[str, Any]:
            out: Dict[str, Any] = {}
            m = re.search(r"['\"]?result['\"]?\s*[:：]\s*['\"]?([^\"',}\n\r]+)", text, re.IGNORECASE)
            if m:
                out["result"] = m.group(1).strip().strip('"\'')
            else:
                m2 = re.search(r"\b(〇|△|×|O|X|o|x)\b", text)
                if m2:
                    out["result"] = m2.group(1)

            m_e = re.search(r"['\"]?evidence['\"]?\s*[:：]\s*([\"'])(.*?)\1", text, re.IGNORECASE | re.DOTALL)
            if m_e:
                out["evidence"] = m_e.group(2).strip()
            else:
                m_e2 = re.search(r"evidence\s*[:：\-]\s*(.+)$", text, re.IGNORECASE | re.DOTALL)
                if m_e2:
                    out["evidence"] = m_e2.group(1).strip()

            m_d = re.search(r"['\"]?details['\"]?\s*[:：]\s*([\"'])(.*?)\1", text, re.IGNORECASE | re.DOTALL)
            if m_d:
                out["details"] = m_d.group(2).strip()
            else:
                m_d2 = re.search(r"details\s*[:：\-]\s*(.+)$", text, re.IGNORECASE | re.DOTALL)
                if m_d2:
                    out["details"] = m_d2.group(1).strip()

            if not out.get("evidence") and text:
                out["evidence"] = text.strip()

            if "result" not in out:
                out["result"] = "△"

            return out

        parsed = _heuristic_parse(text)

    def _build_evidence_list(evidence_field, docs_list: List[Document]):
        evs = []
        if not evidence_field:
            for d in docs_list[:TOP_K]:
                evs.append({"source": d.metadata.get("source"), "excerpt": d.page_content[:400].strip()})
            return evs

        if isinstance(evidence_field, str):
            evs.append({"source": "(model-output)", "excerpt": evidence_field})
            for d in docs_list[:TOP_K]:
                evs.append({"source": d.metadata.get("source"), "excerpt": d.page_content[:300].strip()})
            return evs

        if isinstance(evidence_field, list):
            for item in evidence_field:
                if isinstance(item, dict):
                    src = item.get("source") or item.get("file") or item.get("path") or "(unknown)"
                    exc = item.get("excerpt") or item.get("text") or json.dumps(item, ensure_ascii=False)
                    evs.append({"source": src, "excerpt": exc[:400].strip()})
                else:
                    evs.append({"source": "(model-output)", "excerpt": str(item)[:400]})
            return evs

        evs.append({"source": "(model-output)", "excerpt": str(evidence_field)[:400]})
        return evs

    parsed_evidence = _build_evidence_list(parsed.get("evidence"), docs)
    parsed["evidence_normalized"] = parsed_evidence
    return parsed


def format_b_result(b_result: Dict[str, Any]) -> str:
    """Agent B の構造化結果を日本語の整形テキストに変換する"""
    lines: List[str] = []
    res = b_result.get("result") or b_result.get("status") or "△"
    lines.append(f"判定: {res}")
    
    details = b_result.get("details") or b_result.get("detail") or b_result.get("notes")
    if details:
        lines.append("\n説明:")
        if isinstance(details, str):
            lines.append(details)
        else:
            lines.append(json.dumps(details, ensure_ascii=False, indent=2))

    evs = b_result.get("evidence_normalized") or []
    if evs:
        lines.append("\n根拠 (参照文書と抜粋):")
        for i, e in enumerate(evs, 1):
            src = e.get("source") or "(unknown)"
            excerpt = e.get("excerpt") or ""
            excerpt_clean = excerpt.replace("\n", " ").strip()
            if len(excerpt_clean) > 1200:
                excerpt_clean = excerpt_clean[:1200].rstrip() + " ..."
            lines.append(f"  {i}. source: {src}")
            wrapped = textwrap.fill(excerpt_clean, width=100, subsequent_indent='     ')
            lines.append(textwrap.indent(wrapped, '     '))

    if b_result.get("evidence") and not evs:
        lines.append("\nモデル出力（根拠）:")
        lines.append(str(b_result.get("evidence")))

    return "\n".join(lines)


def find_rule_by_query(rules: List[Dict[str, Any]], query: str) -> Optional[Dict[str, Any]]:
    """ルール一覧から query を元にルールを検索する。ID/パス/タイトル/本文の部分一致で最初のマッチを返す。"""
    q = query.strip().lower()
    for r in rules:
        if r.get("id") and str(r.get("id")).lower() == q:
            return r
    for r in rules:
        if q in (r.get("title") or "").lower() or q in (r.get("path") or "").lower() or q in (r.get("content") or "").lower():
            return r
    return None


def run_cui_mode(rules: List[Dict[str, Any]], vectordb: Chroma):
    """CUI ベースの簡易チャットループ"""
    llm = make_chat_model()

    help_text = (
        "コマンド一覧:\n"
        "  help                      ヘルプ表示\n"
        "  list                      読み込んだルール一覧の一部を表示\n"
        "  show <query>              ルールを表示（id/title の部分一致）\n"
        "  check <query>             指定したルールに対してシステムが従っているか評価（A->B の順）\n"
        "  showfull <summary|b>      直近のチェックで保存された項目の全文表示\n"
        "  ask <自由テキスト>        システム情報に関する RAG 質問\n"
        "  quit                      終了\n"
    )

    print("AgentRAG チャットボット (CUI)。help と入力してください。\n")
    last_store: Dict[str, Any] = {"summary": None, "b": None}

    def print_section(title: str, content: str, max_len: int = 1200):
        """見やすいセクション表示"""
        sep = "=" * 80
        print("\n" + sep)
        print(f"{title}")
        print(sep)
        if content is None:
            print("(なし)\n")
            return

        display_text = content
        if isinstance(content, str) and len(content) > max_len:
            display_text = content[:max_len].rstrip() + "\n...（全文は 'showfull' コマンドで表示可）"

        if isinstance(display_text, str) and display_text.strip().startswith(("{", "[")):
            print(display_text)
        else:
            if isinstance(display_text, str):
                paras = [p.strip() for p in display_text.split("\n\n") if p.strip()]
                for p in paras:
                    wrapped = textwrap.fill(p, width=100)
                    print(wrapped)
                    print()
            else:
                print(str(display_text))
        print(sep + "\n")

    while True:
        try:
            cmd = input("> ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\n終了します。")
            break

        if not cmd:
            continue
        if cmd == "help":
            print(help_text)
            continue
        if cmd == "list":
            for i, r in enumerate(rules[:50], 1):
                print(f"{i}. id={r.get('id')} title={r.get('title')} path={r.get('path')}")
            continue
        if cmd.startswith("show "):
            q = cmd[len("show "):].strip()
            r = find_rule_by_query(rules, q)
            if not r:
                print("ルールが見つかりませんでした。部分文字列で検索してみてください。")
            else:
                print("--- ルール ---")
                print(f"id: {r.get('id')}")
                print(f"title: {r.get('title')}")
                print(f"path: {r.get('path')}")
                print("content:")
                print(r.get("content"))
            continue

        if cmd.startswith("check "):
            q = cmd[len("check "):].strip()
            r = find_rule_by_query(rules, q)
            if not r:
                print("ルールが見つかりません。別のクエリを試してください。\n(例: ルールの一部の語句や id を入力)\n")
                continue

            print(f"選択されたルール: {r.get('title')} (path: {r.get('path')})")

            # Agent A: 要約
            rule_text = r.get("content") or ""
            related = retrieve_related_docs(vectordb, rule_text, k=TOP_K)
            print("[Agent A] ルールと関連ドキュメントから要約を作成しています...")
            summary = agent_a_summarize(llm, rule_text, related)
            last_store["summary"] = summary
            print_section("Agent A - 要約プレビュー", summary)

            # Agent B: 確認
            print("[Agent B] ドキュメントがルールに従っているか評価しています...")
            b_result = agent_b_check(llm, summary, r.get("raw", {}), related)
            last_store["b"] = b_result
            
            b_preview_text = format_b_result(b_result)
            print_section("Agent B - 判定（プレビュー）", b_preview_text)
            
            evs = b_result.get("evidence_normalized") or []
            if evs:
                print("根拠 (参照文書と抜粋):")
                for i, e in enumerate(evs, 1):
                    src = e.get("source") or "(unknown)"
                    excerpt = e.get("excerpt") or ""
                    print(f"  {i}. source: {src}")
                    excerpt_clean = excerpt.replace("\n", " ").strip()
                    if len(excerpt_clean) > 1000:
                        excerpt_clean = excerpt_clean[:1000].rstrip() + " ..."
                    wrapped = textwrap.fill(excerpt_clean, width=100, subsequent_indent='     ')
                    print(textwrap.indent(wrapped, '     '))
                    print()
            else:
                print("(根拠情報はありません)")
            
            result_symbol = b_result.get("result", "△")
            if result_symbol == "〇" or result_symbol == "O" or result_symbol == "o":
                print("補足: 判定は '従っている' と見なされます。必要に応じて関連資料を参照してください。\n")
            elif result_symbol == "×" or result_symbol == "X" or result_symbol == "x":
                print("補足: 判定は '従っていない' です。優先的な対応（修正／設定変更等）が必要です。詳細は関連資料を参照してください。\n")
            else:
                print("補足: 判定は '△'（追加確認が必要）です。関連箇所のログや設定ファイルを追加で提供してください。\n")

            print("評価フローが完了しました。必要に応じて 'showfull summary' や 'showfull b' で全文を表示できます。")
            continue

        if cmd.startswith("showfull "):
            what = cmd[len("showfull "):].strip()
            if what not in ("summary", "b"):
                print("'showfull' の引数は summary|b のいずれかを指定してください。")
                continue
            val = last_store.get(what)
            if val is None:
                print(f"まだ '{what}' の出力がありません。先に 'check <query>' を実行してください。")
                continue
            title_map = {"summary": "Agent A - 要約（全文）", "b": "Agent B - 判定（全文）"}
            
            if isinstance(val, str):
                content = val
            else:
                if what == "b":
                    content = format_b_result(val)
                else:
                    content = json.dumps(val, ensure_ascii=False, indent=2)
            print_section(title_map.get(what, what), content, max_len=10_000)
            continue

        if cmd.startswith("ask "):
            q = cmd[len("ask "):].strip()
            docs = retrieve_related_docs(vectordb, q, k=TOP_K)
            context = "\n\n".join([f"[src:{d.metadata.get('source')}]\n{d.page_content}" for d in docs])
            system = "あなたはシステム情報の検索アシスタントです。ユーザの質問に、関連するドキュメントを参照して簡潔に答えてください。"
            messages = [SystemMessage(content=system), HumanMessage(content=f"質問: {q}\n\n参照文書:\n{context}")]
            resp = llm(messages)
            print(resp.content)
            continue

        if cmd in ("quit", "exit", "q"):
            print("終了します。")
            break

        print("不明なコマンドです。help を表示してください。")


def run_streamlit_mode():
    """Streamlit WebUI モードを実行"""
    try:
        import streamlit as st
    except ImportError:
        print("Streamlit がインストールされていません。pip install streamlit を実行してください。")
        sys.exit(1)
        
    # Streamlit UI の実装
    st.set_page_config(page_title="AgentRAG - Web UI", layout="wide")

    st.markdown("# 🛡️ AgentRAG — セキュリティ規則チェッカー")
    st.markdown("*統一基準対応 RAGチャットボット (Streamlit UI)*")
    st.markdown("---")

    # CSS スタイル
    st.markdown(
        """
        <style>
        * { font-size:13px !important; }
        .stButton>button { padding:4px 8px !important; font-size:13px !important; }
        textarea { font-size:12px !important; }
        
        .main .block-container {
            padding-top: 3rem !important;
            padding-left: 1rem !important;
            padding-right: 1rem !important;
            padding-bottom: 5rem !important;
            max-width: none !important;
            overflow-y: visible !important;
        }
        
        .main h1 {
            margin-top: 0 !important;
            margin-bottom: 0.5rem !important;
            padding-top: 0 !important;
            font-size: 1.8rem !important;
            line-height: 1.2 !important;
        }
        
        .main em {
            font-size: 0.9rem !important;
            color: #666 !important;
            display: block !important;
            margin-bottom: 0.5rem !important;
        }
        
        .main hr {
            margin: 0.5rem 0 1rem 0 !important;
        }
        
        html, body, #root {
            overflow-y: auto !important;
            height: 100% !important;
        }
        
        .main {
            overflow-y: auto !important;
            height: 100vh !important;
            padding-top: 0 !important;
        }
        
        .stSelectbox div[data-baseweb="select"] > div {
            max-height: 300px !important; 
            overflow-y: auto !important;
        }
        
        div[data-baseweb="popover"] {
            max-height: 400px !important;
            overflow-y: auto !important;
        }
        
        header[data-testid="stHeader"] {
            height: 2.5rem !important;
        }
        
        .css-1d391kg {
            padding-top: 1rem !important;
        }
        
        @media (max-width: 768px) {
            .main h1 {
                font-size: 1.5rem !important;
            }
            .main .block-container {
                padding-left: 0.5rem !important;
                padding-right: 0.5rem !important;
            }
        }
        
        @media (max-width: 480px) {
            .main h1 {
                font-size: 1.3rem !important;
            }
            * {
                font-size: 12px !important;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    # キャッシュ付き初期化
    @st.cache_resource
    def get_vectordb():
        docs = load_spec_documents(SPEC_DIR)
        return init_chroma(docs)

    @st.cache_resource
    def get_rules():
        return load_rules_from_dir(RULE_DIR)

    @st.cache_resource
    def get_llm():
        return make_chat_model()

    vectordb = None
    llm = None
    rules = []
    try:
        vectordb = get_vectordb()
        llm = get_llm()
        rules = get_rules()
    except Exception as e:
        st.warning("ベクトルDB や LLM の初期化で警告が出ました。OpenAIキーや依存が正しく設定されているか確認してください。")
        st.exception(e)

    # サイドバー
    with st.sidebar:
        st.markdown("### 📋 メニュー")
        st.write(f"**ルール数**: {len(rules):,}件")
        
        page = st.radio("ページ選択", ["🔍 ルールチェック", "💬 RAG 質問"], label_visibility="collapsed")
        
        st.markdown("---")
        st.markdown("### ⚙️ 設定")
        topk = st.slider("参照ドキュメント数", 1, 10, TOP_K, help="RAG検索で参照するドキュメントの数")
        
        st.markdown("---")
        st.markdown("### 📁 対応ファイル形式")
        st.markdown("""
        **ドキュメント読み込み対応:**
        - 📄 PDF (.pdf)
        - 📝 Word (.docx)
        - 📊 Excel (.xlsx)
        - 📈 PowerPoint (.pptx)
        - 📋 Markdown (.md)
        - 📄 Text (.txt)
        
        **ルール定義:**
        - 📋 JSON (.json)
        """)
        
        st.markdown("---")
        st.caption("💡 `specification/` フォルダにファイルを配置してください")

    # ルール一覧を取得
    def create_rule_preview(rule):
        """ルールの選択肢用プレビューテキストを作成"""
        rule_id = rule.get('id', '')
        title = rule.get('title', '')
        rule_type = rule.get('type', '')
        content = rule.get('content', '')
        
        preview_content = content.replace('\n', ' ').replace('\r', '').strip()
        if len(preview_content) > 50:
            preview_content = preview_content[:50] + "..."
        
        choice_text = f"{rule_id}"
        if rule_type:
            choice_text += f" [{rule_type}]"
        
        if title and title != rule_id:
            short_title = title[:30] + "..." if len(title) > 30 else title
            choice_text += f" {short_title}"
        
        if preview_content:
            choice_text += f" | {preview_content}"
        
        return choice_text

    rule_choices = {}
    for r in rules:
        preview_text = create_rule_preview(r)
        rule_choices[preview_text] = r

    if "ルールチェック" in page:
        st.header("🔍 ルールチェック")
        
        # 検索機能
        col1, col2 = st.columns([3, 1])
        with col1:
            search_term = st.text_input("ルール検索（ID、種別、内容で検索）", placeholder="例: 責任者, 遵守事項, 2.1.1")
        with col2:
            st.write("")
            show_all = st.checkbox("全件表示", help="チェックすると検索結果の全件を表示します（重い場合があります）")
        
        max_display_items = 500 if show_all else 100
        
        # 検索でフィルタリング
        filtered_choices = {}
        
        if search_term:
            count = 0
            for preview_text, rule in rule_choices.items():
                if count >= max_display_items:
                    break
                if (search_term.lower() in preview_text.lower() or
                    search_term.lower() in rule.get('content', '').lower() or
                    search_term.lower() in rule.get('id', '').lower() or
                    search_term.lower() in rule.get('type', '').lower()):
                    filtered_choices[preview_text] = rule
                    count += 1
            
            if filtered_choices:
                total_matches = sum(1 for preview_text, rule in rule_choices.items() 
                                  if (search_term.lower() in preview_text.lower() or
                                      search_term.lower() in rule.get('content', '').lower() or
                                      search_term.lower() in rule.get('id', '').lower() or
                                      search_term.lower() in rule.get('type', '').lower()))
                
                if total_matches > max_display_items:
                    st.info(f"🔍 検索結果: {total_matches}件中 上位{len(filtered_choices)}件を表示")
                    if not show_all:
                        st.caption("より多く表示するには「全件表示」をチェックするか、検索語を具体化してください")
                else:
                    st.success(f"🔍 検索結果: {len(filtered_choices)}件のルールが見つかりました")
            else:
                st.warning("🔍 検索条件に一致するルールが見つかりませんでした")
        else:
            count = 0
            for preview_text, rule in rule_choices.items():
                if count >= max_display_items:
                    break
                filtered_choices[preview_text] = rule
                count += 1
            
            if len(rule_choices) > max_display_items:
                st.info(f"📋 全{len(rule_choices)}件中 上位{max_display_items}件を表示")
                st.caption("検索機能または「全件表示」チェックで他のルールも表示できます")
        
        # ルール選択
        with st.container():
            choices = ["(選択してください)"] + list(filtered_choices.keys())
            sel = st.selectbox(
                "評価するルールを選択", 
                choices, 
                help="ルールID、種別、内容のプレビューが表示されます",
                key="rule_selector"
            )
        
        st.caption("⚙️ ルールを選択して 'チェック実行' を押すと評価が始まります。")

        if sel == "(選択してください)":
            if search_term:
                st.info("上記の検索結果からルールを選択してください")
            else:
                st.info("ルールを選択してください（上部の検索ボックスで絞り込み可能）")
        else:
            r = filtered_choices[sel]
            
            col1, col2 = st.columns([1, 1])
            
            with col1:
                st.markdown(f"**ID**: `{r.get('id')}`")
                if r.get('type'):
                    st.markdown(f"**種別**: {r.get('type')}")
                if r.get('source_file'):
                    st.markdown(f"**ソース**: {r.get('source_file')}")
            
            with col2:
                st.markdown(f"**タイトル**: {r.get('title')}")
                if r.get('path'):
                    st.markdown(f"**階層**: `{r.get('path')}`")
            
            content = r.get('content', '')
            if content:
                st.markdown("**内容:**")
                if len(content) > 500:
                    with st.expander(f"内容を表示（{len(content)}文字）"):
                        st.text(content)
                    st.text(content[:200] + "..." if len(content) > 200 else content)
                else:
                    st.text(content)

            if st.button("チェック実行"):
                try:
                    docs = retrieve_related_docs(vectordb, r.get('content') or r.get('title') or "", k=topk)
                    st.write(f"取得ドキュメント: {len(docs)} チャンク（上位 {topk}）")
                    st.info("要約中...")
                    summary = agent_a_summarize(llm, r.get('content') or '', docs)
                    st.success("要約完了")

                    st.info("評価中...")
                    b_result = agent_b_check(llm, summary, r.get('raw', {}), docs)
                    st.success("評価完了")
                    st.subheader("判定（Agent B）")
                    
                    b_text = format_b_result(b_result)
                    def _normalize_display(text: str) -> str:
                        if not text:
                            return ""
                        t = text.replace('\r\n', '\n').replace('\r', '\n')
                        t = re.sub(r"\n{3,}", "\n\n", t)
                        lines = [ln.rstrip() for ln in t.split('\n')]
                        while lines and lines[0].strip() == "":
                            lines.pop(0)
                        while lines and lines[-1].strip() == "":
                            lines.pop()
                        out_lines = []
                        prev_blank = False
                        for ln in lines:
                            if ln.strip() == "":
                                if not prev_blank:
                                    out_lines.append("")
                                prev_blank = True
                            else:
                                out_lines.append(ln.lstrip())
                                prev_blank = False
                        return "\n".join(out_lines)

                    b_text_clean = _normalize_display(b_text)
                    res_symbol = b_result.get("result") or b_result.get("status") or "△"
                    st.markdown(f"**判定: {res_symbol}**")

                    details = b_result.get("details") or b_result.get("detail") or b_result.get("notes")
                    if details:
                        st.markdown("**説明:**")
                        st.text(details if isinstance(details, str) else json.dumps(details, ensure_ascii=False, indent=2))

                    evs = b_result.get("evidence_normalized") or []
                    if evs:
                        st.markdown("**根拠 (参照文書と抜粋):**")
                        for i, e in enumerate(evs, 1):
                            src = e.get("source") or "(unknown)"
                            excerpt = e.get("excerpt") or ""
                            with st.expander(f"{i}. {src}"):
                                ex = excerpt.replace("\r\n", "\n").replace("\r", "\n").strip()
                                st.text(ex)
                    else:
                        st.info("(根拠情報はありません)")

                    with st.expander("（参考）整形済みテキスト（生）"):
                        st.text(b_text_clean)

                except Exception as e:
                    st.error("評価に失敗しました。ログを確認してください。")
                    st.exception(e)

    elif "RAG" in page:
        st.header("💬 RAG 質問 (システム情報に関する QA)")
        st.caption("📁 PDF, Word, Excel, PowerPoint, Markdown, テキストファイルから情報を検索できます")
        
        q = st.text_input("質問を入力してください", placeholder="例: ウイルス対策の要件は？ / Excel形式の要件は？")
        if st.button("質問実行"):
            if not q:
                st.warning("質問を入力してください")
            else:
                try:
                    docs = retrieve_related_docs(vectordb, q, k=topk)
                    st.write(f"🔍 {len(docs)}件の関連ドキュメントを検索しました")
                    
                    file_types_found = set()
                    for d in docs:
                        file_type = d.metadata.get('file_type', 'unknown')
                        file_types_found.add(file_type)
                    
                    if file_types_found:
                        type_emojis = {'.pdf': '📄', '.docx': '📝', '.xlsx': '📊', '.pptx': '📈', '.md': '📋', '.txt': '📄'}
                        type_str = " ".join([f"{type_emojis.get(ft, '📄')}{ft}" for ft in sorted(file_types_found)])
                        st.caption(f"参照ファイル形式: {type_str}")
                    
                    context = "\n\n".join([f"[src:{d.metadata.get('source')}]\n{d.page_content}" for d in docs])
                    system = "あなたはシステム情報の検索アシスタントです。ユーザの質問に、関連するドキュメントを参照して簡潔に答えてください。\n\n重要: 回答は必ず日本語で行ってください。"
                    messages = [SystemMessage(content=system), HumanMessage(content=f"質問: {q}\n\n参照文書:\n{context}")]
                    resp = llm(messages)
                    st.markdown("**回答:**")
                    st.text(resp.content)
                    
                    with st.expander("🔗 参照元ファイル詳細"):
                        for i, d in enumerate(docs, 1):
                            source = d.metadata.get('source', 'unknown')
                            file_type = d.metadata.get('file_type', 'unknown')
                            chunk_id = d.metadata.get('chunk', 0)
                            st.text(f"{i}. {Path(source).name} ({file_type}, chunk {chunk_id})")
                            st.text(f"   内容: {d.page_content[:100]}...")
                            st.text("")
                            
                except Exception as e:
                    st.error("QA 実行でエラーが発生しました")
                    st.exception(e)

    st.caption("この UI は Streamlit を利用しています。バックエンドである main.py を大きく変更せずにフロントエンドを提供します。")


def main():
    parser = argparse.ArgumentParser(description="セキュリティ規則チェックRAGチャットボット")
    parser.add_argument("--cui", action="store_true", help="CUIモードで実行")
    args = parser.parse_args()
    
    # ルール読み込み
    rules = load_rules_from_dir(RULE_DIR)

    # ドキュメント読み込み
    docs = load_spec_documents(SPEC_DIR)

    # Chroma 初期化・インデックス作成
    vectordb = init_chroma(docs)
    
    if args.cui:
        # CUIモード
        run_cui_mode(rules, vectordb)
    else:
        # StreamlitモードならStreamlitアプリとして実行
        # この場合は本当のエントリーポイントではなく、
        # streamlit run combined_app.py で実行される
        run_streamlit_mode()


if __name__ == "__main__":
    main()